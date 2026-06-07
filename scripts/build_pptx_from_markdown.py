from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


INPUT_MD = Path("docs/sign-language-pipeline-model-presentation-72-slides.md")
OUTPUT_PPTX = Path("docs/sign-language-pipeline-model-presentation-72-slides.pptx")


def split_slides(markdown_text: str) -> list[str]:
    parts = re.split(r"\n---\n", markdown_text.strip())
    return [p.strip() for p in parts if p.strip()]


def parse_slide_content(chunk: str) -> tuple[str, list[str], str | None]:
    lines = [ln.rstrip() for ln in chunk.splitlines()]

    title = ""
    bullets: list[str] = []
    visual_note: str | None = None

    for idx, line in enumerate(lines):
        if line.startswith("# "):
            title = line[2:].strip()
            content_lines = lines[idx + 1 :]
            break
    else:
        title = lines[0].lstrip("# ").strip() if lines else "Slide"
        content_lines = lines[1:] if len(lines) > 1 else []

    in_visual = False
    visual_buffer: list[str] = []

    for raw in content_lines:
        line = raw.strip()
        if not line:
            continue

        if line.lower().startswith("visual:"):
            in_visual = True
            value = line.split(":", 1)[1].strip()
            if value:
                visual_buffer.append(value)
            continue

        if in_visual:
            if line.startswith("-"):
                visual_buffer.append(line[1:].strip())
            else:
                visual_buffer.append(line)
            continue

        if line.startswith("#"):
            bullets.append(line.lstrip("# ").strip())
        elif line.startswith("-"):
            bullets.append(line[1:].strip())
        else:
            bullets.append(line)

    if visual_buffer:
        visual_note = " ".join(visual_buffer)

    return title or "Slide", bullets, visual_note


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(9, 11, 18)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.22))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(46, 196, 244)
    top_bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.7), Inches(0.35), Inches(1.35), Inches(0.32))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(198, 74, 229)
    accent.line.fill.background()


def add_title(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(245, 247, 255)


def add_body(slide, bullets: list[str]):
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.85), Inches(4.65))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets[:9]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 226, 240)
        p.space_after = Pt(10)


def add_visual_note(slide, note: str | None):
    if not note:
        return

    note_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(11.85), Inches(0.95))
    note_card.fill.solid()
    note_card.fill.fore_color.rgb = RGBColor(18, 25, 38)
    note_card.line.color.rgb = RGBColor(46, 196, 244)
    note_card.line.width = Pt(1.5)

    tf = note_card.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Visual direction"
    p1.font.bold = True
    p1.font.size = Pt(13)
    p1.font.color.rgb = RGBColor(111, 224, 255)

    p2 = tf.add_paragraph()
    p2.text = note
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(214, 223, 240)


def add_footer(slide, index: int, total: int):
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.9), Inches(0.3))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Together AI Platform  |  Slide {index}/{total}"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(135, 145, 168)


def build_presentation(slides: list[str]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    total = len(slides)

    for i, chunk in enumerate(slides, start=1):
        title, bullets, visual_note = parse_slide_content(chunk)
        slide = prs.slides.add_slide(blank)

        add_background(slide)
        add_title(slide, title)
        add_body(slide, bullets)
        add_visual_note(slide, visual_note)
        add_footer(slide, i, total)

    prs.save(OUTPUT_PPTX)


def main():
    if not INPUT_MD.exists():
        raise FileNotFoundError(f"Input markdown not found: {INPUT_MD}")

    markdown_text = INPUT_MD.read_text(encoding="utf-8")
    slides = split_slides(markdown_text)
    build_presentation(slides)
    print(f"Generated {OUTPUT_PPTX} with {len(slides)} slides.")


if __name__ == "__main__":
    main()
